"""
PPT Content-Quality Validator  (Agent 1)
=========================================
Implements rules CQ-001 .. CQ-013 from Agent1.xlsx.

Endpoint:  POST /validate
Input (JSON): { "fileName": "...", "region": "AMER|EUR|MENA", "fileContent": "<base64 pptx>" }
Health:    GET /health

Deployed on Vercel (Python serverless runtime).
Vercel loads the top-level `app` (Flask WSGI) as the function handler.
"""

import base64
import binascii
import io
import os
import re

from pptx import Presentation
from pptx.util import Emu


# ------------------------- Region handling ------------------------- #
def normalize_region(region):
    if not region:
        return "AMER/EUR"
    r = region.strip().upper()
    return "MENA" if "MENA" in r else "AMER/EUR"


# ------------------------- Dictionaries ------------------------- #
VAGUE_TERMS = ["improve", "seamless", "cutting-edge", "significant",
               "strong", "better", "various", "many", "several"]
SUPERLATIVES = ["best", "leading", "world-class", "unparalleled",
                "unique", "guaranteed", "proven", "top-tier"]
CASUAL = ["we got", "stuff", "things", "awesome", "quick and dirty",
          "kind of", "sort of", "a lot of"]
PLACEHOLDERS = ["tbd", "tbc", "xxx", "lorem", "[client]", "<insert>", "???"]
APPROVED_ACRONYMS = {"AI", "API", "SAAS", "ROI", "KPI", "SLA", "CRM", "ERP"}
EVIDENCE_RE = re.compile(r"(\d|%|\$|source|client|ref|http)", re.I)
ACRONYM_RE = re.compile(r"\b([A-Z]{2,6})\b")
SENTENCE_SPLIT = re.compile(r"(?<=[.!?])\s+")


# ------------------------- Slide extraction ------------------------- #
def is_footer_shape(shape, slide_height):
    try:
        if shape.top is None or shape.height is None:
            return False
        bottom = shape.top + shape.height
        near_bottom = bottom > slide_height * 0.88
        thin = shape.height <= slide_height * 0.10
        return near_bottom and thin
    except Exception:
        return False


def get_slide_data(prs):
    slides = []
    sh = prs.slide_height or Emu(6858000)
    for idx, slide in enumerate(prs.slides):
        bullets, body_text, footer_text, has_table = [], [], [], False
        for shape in slide.shapes:
            if shape.has_table:
                has_table = True
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            body_text.append(cell.text.strip())
                continue
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            if is_footer_shape(shape, sh):
                footer_text.append(txt)
                continue
            for para in shape.text_frame.paragraphs:
                p = "".join(run.text for run in para.runs).strip()
                if p:
                    bullets.append(p)
            body_text.append(txt)
        slides.append({
            "index": idx + 1, "bullets": bullets,
            "text": "\n".join(body_text), "footer": "\n".join(footer_text),
            "has_table": has_table, "num_shapes": len(slide.shapes),
        })
    return slides


def classify_slide(sd, total):
    words = len(sd["text"].split())
    if sd["index"] == 1:
        return "cover"
    if words <= 8 and len(sd["bullets"]) <= 2:
        return "divider"
    return "content"


# ------------------------- Rule engine ------------------------- #
def add(issues, slide, rid, cap, level, detail, rec):
    issues.append({"slide": slide, "ruleId": rid, "capability": cap,
                   "level": level, "detail": detail, "recommendation": rec})


def count_words(text):
    return len(re.findall(r"\b[\w'-]+\b", text))


def validate_slide(sd, region, kind, issues):
    text = sd["text"]
    words = count_words(text)
    chars = len(re.sub(r"\s", "", text))
    lower = text.lower()

    # CQ-001 Slide density
    if kind == "content":
        t_warn, t_fail = (160, 200) if region == "MENA" else (110, 140)
        if words > t_fail and not sd["has_table"]:
            add(issues, sd["index"], "CQ-001", "Slide density", "fail",
                f"{words} words (fail >{t_fail}).",
                "Split slide, convert prose to bullets, or move detail to appendix.")
        elif words > t_warn and not sd["has_table"]:
            add(issues, sd["index"], "CQ-001", "Slide density", "warn",
                f"{words} words (warn >{t_warn}).", "Trim toward target range.")

    # CQ-002 Character density
    c_fail = 1600 if region == "MENA" else 1200
    if chars > c_fail:
        add(issues, sd["index"], "CQ-002", "Character density", "fail",
            f"{chars} visible characters (fail >{c_fail}).", "Shorten, group, or split.")
    elif chars > 850:
        add(issues, sd["index"], "CQ-002", "Character density", "warn",
            f"{chars} visible characters (warn >850).", "Shorten or group content.")

    # CQ-003 Bullet length
    for b in sd["bullets"]:
        w, c = count_words(b), len(b)
        if w > 25 or c > 160:
            add(issues, sd["index"], "CQ-003", "Bullet length", "fail",
                f'Bullet {w} words / {c} chars: "{b[:60]}..."',
                "Rewrite to a single action/outcome phrase.")
        elif w > 18 or c > 120:
            add(issues, sd["index"], "CQ-003", "Bullet length", "warn",
                f'Bullet {w} words / {c} chars: "{b[:60]}..."', "Tighten toward 6-14 words.")

    # CQ-004 Bullet count
    nb = len(sd["bullets"])
    per_slide_fail = 12 if region == "MENA" else 8
    if kind == "content" and nb > per_slide_fail and not sd["has_table"]:
        add(issues, sd["index"], "CQ-004", "Bullet count", "fail",
            f"{nb} bullets on slide (fail >{per_slide_fail}).",
            "Prioritize top 3-5 or split into two slides.")
    elif nb > 6:
        add(issues, sd["index"], "CQ-004", "Bullet count", "warn",
            f"{nb} bullets (target 3-5).", "Merge or prioritize.")

    # CQ-005 Sentence length
    for s in SENTENCE_SPLIT.split(text):
        w = count_words(s)
        if w > 35:
            add(issues, sd["index"], "CQ-005", "Sentence length", "fail",
                f'{w}-word sentence: "{s[:60]}..."',
                "Break into two sentences or convert to a bullet.")
        elif w > 28:
            add(issues, sd["index"], "CQ-005", "Sentence length", "warn",
                f'{w}-word sentence: "{s[:60]}..."', "Shorten toward 12-22 words.")

    # CQ-006 Paragraph length
    for para in text.split("\n"):
        w = count_words(para)
        if w > 70:
            add(issues, sd["index"], "CQ-006", "Paragraph length", "fail",
                f"{w}-word paragraph.", "Convert to bullets or summary + detail.")
        elif w > 50:
            add(issues, sd["index"], "CQ-006", "Paragraph length", "warn",
                f"{w}-word paragraph.", "Break into bullets.")

    # CQ-007 Weak/vague wording
    vague_hits = [t for t in VAGUE_TERMS if re.search(rf"\b{re.escape(t)}\b", lower)
                  and not EVIDENCE_RE.search(text)]
    if len(vague_hits) > 3:
        add(issues, sd["index"], "CQ-007", "Weak/vague wording", "fail",
            f"Vague terms without evidence: {', '.join(vague_hits)}.",
            "Replace with measurable verbs/outcomes.")
    elif vague_hits:
        add(issues, sd["index"], "CQ-007", "Weak/vague wording", "warn",
            f"Vague term(s): {', '.join(vague_hits)}.",
            "Add a metric or replace with a concrete outcome.")

    # CQ-008 Unsupported superlatives
    sup_hits = [t for t in SUPERLATIVES if re.search(rf"\b{re.escape(t)}\b", lower)]
    if sup_hits and not EVIDENCE_RE.search(text):
        add(issues, sd["index"], "CQ-008", "Unsupported superlatives", "warn",
            f"Superlative(s) without evidence: {', '.join(sup_hits)}.",
            "Soften the claim or add a source/number.")

    # CQ-009 Professional tone
    casual_hits = [t for t in CASUAL if t in lower]
    excl = "!" in text and kind not in ("cover",)
    if casual_hits or excl:
        d = []
        if casual_hits:
            d.append("casual: " + ", ".join(casual_hits))
        if excl:
            d.append("exclamation mark used")
        add(issues, sd["index"], "CQ-009", "Professional tone", "warn",
            "; ".join(d) + ".", "Rewrite in a formal, professional tone.")

    # CQ-011 Completeness
    ph = [p for p in PLACEHOLDERS if p in lower]
    if ph:
        add(issues, sd["index"], "CQ-011", "Completeness", "fail",
            f"Placeholder(s) present: {', '.join(ph)}.",
            "Provide the missing content before sharing.")

    # CQ-012 Acronym hygiene
    for ac in set(ACRONYM_RE.findall(text)):
        if ac.upper() in APPROVED_ACRONYMS:
            continue
        if re.search(rf"\([^)]*{ac}[^)]*\)", text) or re.search(rf"{ac}\s*\(", text):
            continue
        add(issues, sd["index"], "CQ-012", "Acronym hygiene", "warn",
            f'Acronym "{ac}" not expanded on first use.',
            "Spell out the acronym the first time it appears.")

    # CQ-013 Number formatting
    fmt = []
    if re.search(r"\d+\s+[Xx]\s+\d+", text):
        fmt.append("use 24x7 not 24 X 7")
    if re.search(r"\d\s%", text):
        fmt.append("no space before %")
    if re.search(r"\$\s*\d+\s*(trillion|billion|million)", text, re.I):
        fmt.append("use compact $6T/$6B/$6M")
    if fmt:
        add(issues, sd["index"], "CQ-013", "Number formatting", "warn",
            "; ".join(fmt) + ".", "Normalize metrics to the style guide.")


def validate_presentation(prs, region):
    slides = get_slide_data(prs)
    total = len(slides)
    issues = []
    for sd in slides:
        validate_slide(sd, region, classify_slide(sd, total), issues)

    counts = {"fail": 0, "warn": 0}
    for i in issues:
        counts[i["level"]] = counts.get(i["level"], 0) + 1

    if counts["fail"]:
        verdict = f"NEEDS WORK — {counts['fail']} failing and {counts['warn']} warning issue(s)."
    elif counts["warn"]:
        verdict = f"MOSTLY OK — {counts['warn']} warning issue(s), no failures."
    else:
        verdict = "PASS — no content-quality issues detected."

    return {"summary": verdict, "region": region, "totalSlides": total,
            "counts": counts, "issues": issues}


# ------------------------- Base64 -> Presentation ------------------------- #
def load_presentation_from_base64(b64):
    if b64 is None:
        raise ValueError("fileContent (base64) is missing.")
    s = b64.strip()
    if s.startswith("data:") and "," in s:
        s = s.split(",", 1)[1]
    s = re.sub(r"[^A-Za-z0-9+/=]", "", s)
    s += "=" * (-len(s) % 4)
    try:
        raw = base64.b64decode(s)
    except (binascii.Error, ValueError) as e:
        raise ValueError(f"Base64 decode failed: {e}")
    if raw[:2] != b"PK":
        raise ValueError("Decoded content is not a valid .pptx (missing PK zip header).")
    return Presentation(io.BytesIO(raw))


def run_validation(payload):
    region = normalize_region(payload.get("region"))
    prs = load_presentation_from_base64(payload.get("fileContent")
                                        or payload.get("$content")
                                        or payload.get("contentBytes"))
    result = validate_presentation(prs, region)
    result["fileName"] = payload.get("fileName", "presentation.pptx")
    return result


# ------------------------- Flask app ------------------------- #
from flask import Flask, request, jsonify
app = Flask(__name__)


@app.route("/validate", methods=["POST"])
def validate():
    try:
        payload = request.get_json(force=True, silent=True) or {}
        return jsonify(run_validation(payload)), 200
    except Exception as e:
        return jsonify({"error": str(e), "summary": "Validation failed.",
                        "issues": []}), 200


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"}), 200


@app.route("/", methods=["GET"])
def home():
    return jsonify({"service": "PPT Content Quality Validator", "status": "running"}), 200


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
